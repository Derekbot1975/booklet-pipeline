"""
Reference Document Library for the Booklet Pipeline.

Manages the school's reference documents that inform AI generation:
  - Expert input files (per subject)
  - National Curriculum documents (per subject + key stage)
  - KS4/5 specifications (per exam board + subject)
  - Scheme of Work exemplar format
  - Great Lesson / Great Booklet codification

Documents are stored in data/reference-docs/{id}/ with:
  - metadata.json  — category, subject, key_stage, exam_board, title, etc.
  - original file   — the uploaded .docx, .pdf, or .xlsx
  - parsed.txt      — AI-extracted plain text content

The parsed content is injected into Claude prompts as context for quality assurance.
"""

import hashlib
import json
import logging
import os
import re
import time
import uuid
from datetime import datetime
from pathlib import Path

logger = logging.getLogger(__name__)

DOCS_DIR = Path(__file__).parent / "data" / "reference-docs"
DOCS_DIR.mkdir(parents=True, exist_ok=True)

# Valid categories
CATEGORIES = {
    "expert_input": "Expert Input Files",
    "national_curriculum": "National Curriculum",
    "ks4_spec": "KS4/5 Specifications",
    "sow_exemplar": "Scheme of Work Exemplar Format",
    "great_lesson_code": "Great Lesson / Great Booklet Codification",
}

VALID_EXTENSIONS = {".docx", ".pdf", ".xlsx", ".xls", ".txt", ".md", ".pptx", ".pages"}


def _doc_dir(doc_id):
    """Get the directory for a specific document."""
    return DOCS_DIR / doc_id


def _meta_path(doc_id):
    """Get the metadata file path for a document."""
    return _doc_dir(doc_id) / "metadata.json"


def _parsed_path(doc_id):
    """Get the parsed content file path for a document."""
    return _doc_dir(doc_id) / "parsed.txt"


def _file_hash(filepath):
    """Compute SHA-256 hash of a file for change detection."""
    h = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def list_documents(category=None, subject=None, key_stage=None):
    """List all reference documents, optionally filtered."""
    docs = []
    for d in sorted(DOCS_DIR.iterdir()):
        meta_file = d / "metadata.json"
        if not meta_file.exists():
            continue
        try:
            meta = json.loads(meta_file.read_text())
        except (json.JSONDecodeError, OSError):
            continue

        if category and meta.get("category") != category:
            continue
        if subject and meta.get("subject") != subject:
            continue
        if key_stage and meta.get("key_stage") != key_stage:
            continue

        # Add parsed status
        meta["has_parsed_content"] = _parsed_path(meta["id"]).exists()
        docs.append(meta)

    return docs


def get_document(doc_id):
    """Get a single document's metadata and parsed content."""
    meta_file = _meta_path(doc_id)
    if not meta_file.exists():
        return None

    meta = json.loads(meta_file.read_text())
    parsed_file = _parsed_path(doc_id)
    if parsed_file.exists():
        meta["parsed_content"] = parsed_file.read_text()
    else:
        meta["parsed_content"] = None
    meta["has_parsed_content"] = parsed_file.exists()
    return meta


def save_document(file_path, category, title, subject=None, key_stage=None,
                  exam_board=None, description=None, doc_id=None):
    """
    Save a reference document. Copies the file, stores metadata,
    and parses content.

    Args:
        file_path: Path to the uploaded file
        category: One of the CATEGORIES keys
        title: Human-readable title
        subject: Optional subject (e.g. "Science", "History")
        key_stage: Optional key stage (e.g. "KS3", "KS4")
        exam_board: Optional exam board (e.g. "AQA", "Edexcel")
        description: Optional description
        doc_id: Optional existing ID (for replacement)

    Returns:
        Document metadata dict
    """
    if category not in CATEGORIES:
        raise ValueError(f"Invalid category: {category}. Must be one of: {list(CATEGORIES.keys())}")

    src = Path(file_path)
    if not src.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = src.suffix.lower()
    if ext not in VALID_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {VALID_EXTENSIONS}")

    # Generate or reuse ID
    if not doc_id:
        doc_id = str(uuid.uuid4())[:8]

    doc_dir = _doc_dir(doc_id)
    doc_dir.mkdir(parents=True, exist_ok=True)

    # Copy file to doc directory
    import shutil
    dest_filename = f"document{ext}"
    dest_path = doc_dir / dest_filename
    shutil.copy2(str(src), str(dest_path))

    # Build metadata
    content_hash = _file_hash(str(dest_path))
    now = datetime.utcnow().isoformat()

    meta = {
        "id": doc_id,
        "category": category,
        "category_label": CATEGORIES[category],
        "subject": subject,
        "key_stage": key_stage,
        "exam_board": exam_board,
        "title": title,
        "description": description or "",
        "filename": src.name,
        "file_format": ext.lstrip("."),
        "file_path": str(dest_path),
        "content_hash": content_hash,
        "uploaded_at": now,
        "updated_at": now,
    }

    _meta_path(doc_id).write_text(json.dumps(meta, indent=2))

    # Parse the document content
    try:
        parsed = _parse_document(str(dest_path), ext)
        _parsed_path(doc_id).write_text(parsed)
        meta["has_parsed_content"] = True
        logger.info(f"Parsed reference doc '{title}' ({len(parsed)} chars)")
    except Exception as e:
        logger.error(f"Failed to parse reference doc '{title}': {e}")
        meta["has_parsed_content"] = False

    return meta


def update_document(doc_id, **kwargs):
    """Update metadata fields for an existing document."""
    meta_file = _meta_path(doc_id)
    if not meta_file.exists():
        raise ValueError(f"Document not found: {doc_id}")

    meta = json.loads(meta_file.read_text())

    allowed_fields = {"title", "description", "subject", "key_stage", "exam_board", "category"}
    for key, value in kwargs.items():
        if key in allowed_fields:
            meta[key] = value
            if key == "category" and value in CATEGORIES:
                meta["category_label"] = CATEGORIES[value]

    meta["updated_at"] = datetime.utcnow().isoformat()
    meta_file.write_text(json.dumps(meta, indent=2))
    return meta


def delete_document(doc_id):
    """Delete a reference document and all its files."""
    doc_dir = _doc_dir(doc_id)
    if not doc_dir.exists():
        raise ValueError(f"Document not found: {doc_id}")

    import shutil
    shutil.rmtree(str(doc_dir))
    return True


def reparse_document(doc_id):
    """Re-parse a document's content (e.g. after updating the parser)."""
    meta = get_document(doc_id)
    if not meta:
        raise ValueError(f"Document not found: {doc_id}")

    file_path = meta.get("file_path")
    if not file_path or not Path(file_path).exists():
        raise FileNotFoundError(f"Source file missing for document {doc_id}")

    ext = Path(file_path).suffix.lower()
    parsed = _parse_document(file_path, ext)
    _parsed_path(doc_id).write_text(parsed)
    return {"doc_id": doc_id, "parsed_length": len(parsed)}


def get_reference_context(subject=None, key_stage=None, exam_board=None,
                          max_chars=480_000):
    """
    Build the reference context string for AI prompts.

    Gathers all relevant reference documents and assembles them into a
    single context string that gets injected into Claude system prompts.
    Uses a character budget (~120K tokens at ~4 chars/token) to stay well
    within API limits.

    Priority order:
      1. great_lesson_code  — always included (small, cross-cutting)
      2. Matching spec       — the specific exam board specification
      3. National curriculum — matching subject + key stage
      4. Expert input        — matching subject, individually capped
      5. sow_exemplar        — format examples (large, lower priority)

    Per-document caps prevent any single huge PDF from consuming the budget.
    """
    MAX_SPEC_CHARS = 60_000       # ~15K tokens per spec
    MAX_EXPERT_CHARS = 20_000     # ~5K tokens per expert doc
    MAX_NC_CHARS = 40_000         # ~10K tokens per NC doc
    MAX_EXEMPLAR_CHARS = 40_000   # ~10K tokens per exemplar

    docs = list_documents()

    # ── Bucket docs by category ──────────────────────────────────
    buckets = {
        "great_lesson_code": [],
        "ks4_spec": [],
        "national_curriculum": [],
        "expert_input": [],
        "sow_exemplar": [],
    }

    for doc in docs:
        doc_id = doc["id"]
        parsed_file = _parsed_path(doc_id)
        if not parsed_file.exists():
            continue
        cat = doc.get("category")
        if cat not in buckets:
            continue

        # Apply category-specific subject/filter matching
        if cat == "great_lesson_code":
            buckets[cat].append(doc)
        elif cat == "sow_exemplar":
            buckets[cat].append(doc)
        elif cat == "expert_input":
            if subject and (doc.get("subject") or "").lower() == subject.lower():
                buckets[cat].append(doc)
        elif cat == "national_curriculum":
            doc_subj = (doc.get("subject") or "").lower()
            doc_ks = (doc.get("key_stage") or "").upper()
            if subject and doc_subj == subject.lower():
                if not key_stage or doc_ks == key_stage.upper():
                    buckets[cat].append(doc)
        elif cat == "ks4_spec":
            doc_subj = (doc.get("subject") or "").lower()
            doc_board = (doc.get("exam_board") or "").lower()
            if subject and doc_subj == subject.lower():
                if not exam_board or doc_board == exam_board.lower():
                    buckets[cat].append(doc)

    # ── Assemble context with budget tracking ────────────────────
    context_parts = []
    budget = max_chars

    def _add(header, title, content, cap):
        nonlocal budget
        if budget <= 0:
            return False
        if len(content) > cap:
            content = content[:cap] + "\n\n[... truncated for length ...]"
        entry = f"{header}\nTitle: {title}\n\n{content}"
        if len(entry) > budget:
            entry = entry[:budget]
        context_parts.append(entry)
        budget -= len(entry)
        return True

    # 1. Great lesson code (always, small docs)
    for doc in buckets["great_lesson_code"]:
        content = _parsed_path(doc["id"]).read_text().strip()
        if not content:
            continue
        _add("=== SCHOOL QUALITY FRAMEWORK ===",
             doc.get("title", "Quality Framework"), content, 30_000)

    # 2. Matching specification(s)
    for doc in buckets["ks4_spec"]:
        content = _parsed_path(doc["id"]).read_text().strip()
        if not content:
            continue
        board_label = doc.get("exam_board") or "Unknown"
        _add(f"=== SPECIFICATION ({board_label} {subject}) ===",
             doc.get("title", "Specification"), content, MAX_SPEC_CHARS)

    # 3. National curriculum
    for doc in buckets["national_curriculum"]:
        content = _parsed_path(doc["id"]).read_text().strip()
        if not content:
            continue
        ks_label = (doc.get("key_stage") or "").upper()
        _add(f"=== NATIONAL CURRICULUM ({ks_label} {subject}) ===",
             doc.get("title", "National Curriculum"), content, MAX_NC_CHARS)

    # 4. Expert input (most numerous — sorted smallest first to fit more)
    expert_docs = []
    for doc in buckets["expert_input"]:
        p = _parsed_path(doc["id"])
        content = p.read_text().strip()
        if content:
            expert_docs.append((len(content), doc, content))
    expert_docs.sort(key=lambda x: x[0])  # smallest first → maximise coverage

    for _, doc, content in expert_docs:
        if budget <= 0:
            break
        _add(f"=== EXPERT INPUT FOR {subject.upper()} ===",
             doc.get("title", "Expert Input"), content, MAX_EXPERT_CHARS)

    # 5. SoW exemplars (large, lower priority — only if budget remains)
    for doc in buckets["sow_exemplar"]:
        if budget <= 0:
            break
        content = _parsed_path(doc["id"]).read_text().strip()
        if not content:
            continue
        _add("=== SCHEME OF WORK FORMAT ===",
             doc.get("title", "SoW Exemplar"), content, MAX_EXEMPLAR_CHARS)

    if not context_parts:
        return ""

    return "\n\n---\n\n".join(context_parts)


# ───────────────────────────────────────────────────────────────────
# Document parsing — extract text from uploaded files
# ───────────────────────────────────────────────────────────────────

def _parse_document(file_path, ext):
    """Parse a document to extract its text content."""
    ext = ext.lower()
    if ext in (".txt", ".md"):
        return Path(file_path).read_text(encoding="utf-8")
    elif ext == ".docx":
        return _parse_docx(file_path)
    elif ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext in (".xlsx", ".xls"):
        return _parse_excel(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    elif ext == ".pages":
        return _parse_pages(file_path)
    else:
        raise ValueError(f"Cannot parse file type: {ext}")


def _parse_docx(file_path):
    """Extract text from a .docx file."""
    try:
        from docx import Document
        doc = Document(file_path)
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Preserve heading structure
                if para.style and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "").strip()
                    try:
                        level_num = int(level)
                        parts.append(f"{'#' * level_num} {text}")
                    except ValueError:
                        parts.append(text)
                else:
                    parts.append(text)

        # Also extract table content
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))

        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"DOCX parse failed: {e}")
        return f"[Parse error: {e}]"


def _parse_pdf(file_path):
    """Extract text from a PDF file using available libraries."""
    # Try PyPDF2 first (most common)
    try:
        import PyPDF2
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    parts.append(text.strip())
            return "\n\n---\n\n".join(parts)
    except ImportError:
        pass

    # Try pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            parts = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    parts.append(text.strip())
            return "\n\n---\n\n".join(parts)
    except ImportError:
        pass

    # Fallback: try to use subprocess with pdftotext
    try:
        import subprocess
        result = subprocess.run(
            ["pdftotext", "-layout", file_path, "-"],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    return "[PDF parsing requires PyPDF2 or pdfplumber. Install with: pip install PyPDF2]"


def _parse_pptx(file_path):
    """Extract text from a PowerPoint .pptx file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            slide_texts.append(" | ".join(cells))
            if slide_texts:
                parts.append(f"## Slide {i}\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except ImportError:
        logger.error("python-pptx not installed. Run: pip install python-pptx")
        return "[PPTX parse error: python-pptx not installed]"
    except Exception as e:
        logger.error(f"PPTX parse failed: {e}")
        return f"[PPTX parse error: {e}]"


def _parse_pages(file_path):
    """Extract text from an Apple Pages file (which is a zip containing IWA/XML)."""
    import zipfile
    try:
        with zipfile.ZipFile(file_path, 'r') as z:
            # Pages files are zips. Try to find readable text content.
            # The main content is in Index/Document.iwa (protobuf) — hard to parse.
            # Fallback: look for any preview PDF inside the zip.
            names = z.namelist()

            # Try preview PDF first
            pdf_names = [n for n in names if n.lower().endswith('.pdf')]
            if pdf_names:
                import tempfile
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                    tmp.write(z.read(pdf_names[0]))
                    tmp_path = tmp.name
                result = _parse_pdf(tmp_path)
                os.unlink(tmp_path)
                return result

            # Try plain text extraction from buildVersionHistory
            text_parts = []
            for name in names:
                if name.endswith('.txt') or name.endswith('.xml'):
                    try:
                        content = z.read(name).decode('utf-8', errors='ignore')
                        if content.strip():
                            text_parts.append(content.strip())
                    except Exception:
                        pass
            if text_parts:
                return "\n\n".join(text_parts)

            return "[Pages file: could not extract text content. Consider exporting as PDF or DOCX.]"
    except Exception as e:
        logger.error(f"Pages parse failed: {e}")
        return f"[Pages parse error: {e}]"


def _parse_excel(file_path):
    """Extract text from an Excel spreadsheet."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            parts.append("\n".join(rows))
        wb.close()
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"Excel parse failed: {e}")
        return f"[Excel parse error: {e}]"
