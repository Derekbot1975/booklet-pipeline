"""
Lesson Presentation Generator (Prompt Sheet 16).

Generates PowerPoint presentations from lesson data using:
  - Claude API to generate slide content and speaker notes
  - python-pptx to render as .pptx files

Supports individual lessons, units, and full terms.
"""

import json
import logging
import os
import re
import time
import uuid
from datetime import datetime
from pathlib import Path

import anthropic
from dotenv import load_dotenv

load_dotenv(override=True)

logger = logging.getLogger(__name__)


def _extract_json(raw, expect_object=False):
    """Robustly extract JSON from AI response text."""
    raw = re.sub(r"^```(?:json)?\s*\n?", "", raw)
    raw = re.sub(r"\n?```\s*$", "", raw)
    if expect_object:
        start, end = raw.find("{"), raw.rfind("}")
    else:
        start, end = raw.find("["), raw.rfind("]")
    if start != -1 and end != -1 and end > start:
        raw = raw[start:end + 1]
    raw = re.sub(r",\s*([}\]])", r"\1", raw)
    return raw


PRESENTATIONS_DIR = Path(__file__).parent / "data" / "presentations"
PRESENTATIONS_DIR.mkdir(parents=True, exist_ok=True)


SLIDE_GENERATION_PROMPT = """You are a specialist educational presentation designer creating lesson slides.

LANGUAGE: Use UK English throughout.

SLIDE DESIGN PRINCIPLES:
- One key concept per slide
- Minimal text: max 6 bullet points, max 8 words each
- Every slide should have a visual element described in [IMAGE: description]
- Speaker notes contain the FULL teaching script
- High contrast, large font (indicate important text with CAPS or emphasis)
- Variety in slide layouts

OUTPUT FORMAT: You MUST return ONLY a valid JSON array. No markdown, no code fences, no commentary — just the raw JSON array starting with [ and ending with ].

Each slide object in the array must have these keys:
{
    "slideNumber": 1,
    "type": "title",
    "title": "Slide title",
    "content": ["Bullet 1", "Bullet 2"],
    "imageDescription": "Description of visual for this slide",
    "speakerNotes": "Full teaching script for this slide.",
    "layout": "title_slide"
}

Valid "type" values: title, starter, vocabulary, content, worked_example, practice, misconception, plenary, reflection
Valid "layout" values: title_slide, content_left_image_right, full_content, two_column, blank_with_title
"imageDescription" can be null if no image is needed.
"speakerNotes" should include timing, questions to ask, expected answers, and differentiation tips."""


def _get_client():
    return anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))


# ───────────────────────────────────────────────────────────────────
# Slide content generation via Claude
# ───────────────────────────────────────────────────────────────────

def generate_slide_content(lesson, reference_context="",
                           existing_booklet_content=None,
                           include_speaker_notes=True,
                           include_differentiation=True,
                           model="claude-sonnet-4-5-20250929"):
    """Generate slide content JSON for a single lesson."""
    client = _get_client()
    start = time.time()

    user_parts = [
        f"Generate a complete lesson presentation (15-20 slides) for:",
        f"  Subject: {lesson.get('subject', 'Unknown')}",
        f"  Topic: {lesson.get('topic', 'Unknown')}",
        f"  Year: {lesson.get('year', '?')}",
        f"  Lesson Number: {lesson.get('lesson_number', '?')}",
        f"  Title: {lesson.get('title', 'Unknown')}",
        f"  Spec Content: {lesson.get('spec_content', 'N/A')}",
        f"  Key Vocabulary: {lesson.get('key_vocabulary', 'N/A')}",
    ]

    if lesson.get("rp"):
        user_parts.append(f"  Required Practical: {lesson['rp']}")

    user_parts.append(f"\nSLIDE SEQUENCE:")
    user_parts.append(f"1. Title slide (lesson title, date placeholder, objectives)")
    user_parts.append(f"2. Do Now / Starter (3-5 retrieval questions, 5 min)")
    user_parts.append(f"3-4. Key Vocabulary slides")
    user_parts.append(f"5-10. Core Content slides (one concept each)")
    user_parts.append(f"11-12. Worked Example (step by step)")
    user_parts.append(f"13-15. Practice Questions (scaffolded: bronze/silver/gold)")
    user_parts.append(f"16. Misconception Check (true/false or spot the error)")
    user_parts.append(f"17. Plenary (summary + exit ticket)")
    user_parts.append(f"18. Reflection (confidence rating)")

    if existing_booklet_content:
        user_parts.append(f"\nEXISTING BOOKLET CONTENT (use as source):\n{existing_booklet_content[:4000]}")

    if reference_context:
        user_parts.append(f"\nREFERENCE MATERIALS:\n{reference_context[:3000]}")

    if include_speaker_notes:
        user_parts.append(f"\nInclude detailed speaker notes with teaching script, timing, and questioning prompts.")
    if include_differentiation:
        user_parts.append(f"\nInclude differentiation in speaker notes (higher/lower/SEND).")

    user_parts.append(f"\nReturn a JSON array of slide objects.")

    response = client.messages.create(
        model=model,
        max_tokens=8000,
        system=[{
            "type": "text",
            "text": SLIDE_GENERATION_PROMPT,
            "cache_control": {"type": "ephemeral"},
        }],
        messages=[{"role": "user", "content": "\n".join(user_parts)}],
    )

    duration = round(time.time() - start, 1)
    raw = response.content[0].text.strip()

    raw = _extract_json(raw, expect_object=False)  # slides are an array

    try:
        slides = json.loads(raw)
    except json.JSONDecodeError as e:
        logger.error(f"Failed to parse slides JSON: {e}\nRaw (first 800): {raw[:800]}")
        raise ValueError(f"AI generated invalid JSON for slides: {e}")

    if not isinstance(slides, list):
        slides = [slides]

    return {
        "slides": slides,
        "usage": {
            "input_tokens": response.usage.input_tokens,
            "output_tokens": response.usage.output_tokens,
        },
        "model": model,
        "duration_s": duration,
    }


# ───────────────────────────────────────────────────────────────────
# PowerPoint rendering
# ───────────────────────────────────────────────────────────────────

def render_pptx(slides, lesson, output_path=None, style="clean"):
    """
    Render slide data into a PowerPoint file using python-pptx.

    Args:
        slides: List of slide dicts from generate_slide_content()
        lesson: Lesson data dict
        output_path: Optional output file path
        style: "clean" | "subject_themed" | "bold"

    Returns:
        Path to the generated .pptx file
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Subject-based colour theming
    subject = (lesson.get("subject") or "General").lower()
    colour_map = {
        "biology": {"primary": RGBColor(0x22, 0xC5, 0x5E), "secondary": RGBColor(0x16, 0xA3, 0x4A)},
        "chemistry": {"primary": RGBColor(0x3B, 0x82, 0xF6), "secondary": RGBColor(0x25, 0x63, 0xEB)},
        "physics": {"primary": RGBColor(0xF5, 0x9E, 0x0B), "secondary": RGBColor(0xD9, 0x77, 0x06)},
        "science": {"primary": RGBColor(0x8B, 0x5C, 0xF6), "secondary": RGBColor(0x7C, 0x3A, 0xED)},
        "history": {"primary": RGBColor(0xEF, 0x44, 0x44), "secondary": RGBColor(0xDC, 0x26, 0x26)},
        "geography": {"primary": RGBColor(0x06, 0xB6, 0xD4), "secondary": RGBColor(0x08, 0x91, 0xB2)},
        "maths": {"primary": RGBColor(0xF9, 0x73, 0x16), "secondary": RGBColor(0xEA, 0x58, 0x0C)},
        "english": {"primary": RGBColor(0xEC, 0x48, 0x99), "secondary": RGBColor(0xDB, 0x27, 0x77)},
    }
    colours = colour_map.get(subject, {"primary": RGBColor(0x00, 0x7A, 0xFF), "secondary": RGBColor(0x00, 0x66, 0xD6)})

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")
        title_text = slide_data.get("title", "")
        content_items = slide_data.get("content", [])
        speaker_notes = slide_data.get("speakerNotes", "")
        image_desc = slide_data.get("imageDescription", "")

        # Use blank layout and add shapes manually for full control
        blank_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(blank_layout)

        if slide_type == "title":
            # Title slide — large centred text
            _add_title_slide(slide, prs, title_text, content_items, colours, lesson)
        else:
            # Content slide — title bar + content area
            _add_content_slide(slide, prs, title_text, content_items,
                               image_desc, colours, slide_type)

        # Speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    # Save
    if not output_path:
        from generator import OUTPUT_DIR, check_existing_booklet
        existing = check_existing_booklet(lesson)
        if existing["exists"]:
            base_dir = Path(existing["docx_path"]).parent
            base_stem = Path(existing["docx_path"]).stem
        else:
            base_dir = OUTPUT_DIR / "presentations"
            ln = lesson.get("lesson_number", 0)
            title = lesson.get("title", "Untitled")
            safe_title = re.sub(r"[^\w\s-]", "", title).replace(" ", "_")
            base_stem = f"L{ln:03d}_-_{safe_title}"

        base_dir.mkdir(parents=True, exist_ok=True)
        output_path = str(base_dir / f"{base_stem} - Presentation.pptx")

    prs.save(output_path)
    return output_path


def _add_title_slide(slide, prs, title_text, content_items, colours, lesson):
    """Add a title slide with large centred text."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    # Background colour bar at top
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colours["primary"]
    shape.line.fill.background()

    # Lesson title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text or lesson.get("title", "Lesson")
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle / objectives
    if content_items:
        for item in content_items[:4]:
            p2 = tf.add_paragraph()
            p2.text = str(item)
            p2.font.size = Pt(20)
            p2.alignment = PP_ALIGN.CENTER
            from pptx.dml.color import RGBColor as RC
            p2.font.color.rgb = RC(0x66, 0x66, 0x66)

    # Subject and year badge
    meta = f"{lesson.get('subject', '')} — Year {lesson.get('year', '')} — Lesson {lesson.get('lesson_number', '')}"
    txMeta = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.5))
    pMeta = txMeta.text_frame.paragraphs[0]
    pMeta.text = meta
    pMeta.font.size = Pt(16)
    pMeta.alignment = PP_ALIGN.CENTER
    from pptx.dml.color import RGBColor as RC
    pMeta.font.color.rgb = RC(0x99, 0x99, 0x99)


def _add_content_slide(slide, prs, title_text, content_items, image_desc, colours, slide_type):
    """Add a standard content slide."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    # Title bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = colours["primary"]
    bar.line.fill.background()

    txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.7))
    p = txTitle.text_frame.paragraphs[0]
    p.text = title_text or ""
    p.font.size = Pt(28)
    p.font.bold = True
    from pptx.dml.color import RGBColor as RC
    p.font.color.rgb = RC(0xFF, 0xFF, 0xFF)

    # Content area
    content_width = Inches(8) if image_desc else Inches(12)
    txContent = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), content_width, Inches(5.5)
    )
    tf = txContent.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items or []):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}" if not str(item).startswith("•") else str(item)
        p.font.size = Pt(22)
        p.space_after = Pt(8)

    # Image placeholder
    if image_desc:
        img_box = slide.shapes.add_shape(
            1, Inches(9), Inches(1.3), Inches(3.8), Inches(4)
        )
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = RC(0xF5, 0xF5, 0xF7)
        from pptx.util import Emu
        img_box.line.color.rgb = RC(0xD2, 0xD2, 0xD7)

        txImg = slide.shapes.add_textbox(Inches(9.2), Inches(2.5), Inches(3.4), Inches(2))
        pImg = txImg.text_frame.paragraphs[0]
        pImg.text = f"[IMAGE]\n{image_desc}"
        pImg.font.size = Pt(12)
        pImg.font.color.rgb = RC(0x86, 0x86, 0x8B)
        pImg.alignment = PP_ALIGN.CENTER


# ───────────────────────────────────────────────────────────────────
# Presentation CRUD
# ───────────────────────────────────────────────────────────────────

def save_presentation_data(lesson, slides, pptx_path, usage, model, duration_s):
    """Save presentation metadata to data/presentations/."""
    pres_id = str(uuid.uuid4())[:8]
    data = {
        "id": pres_id,
        "lesson_year": lesson.get("year"),
        "lesson_number": lesson.get("lesson_number"),
        "lesson_title": lesson.get("title"),
        "subject": lesson.get("subject"),
        "topic": lesson.get("topic"),
        "pptx_path": pptx_path,
        "slide_count": len(slides),
        "usage": usage,
        "model": model,
        "duration_s": duration_s,
        "created_at": datetime.utcnow().isoformat(),
    }
    path = PRESENTATIONS_DIR / f"{pres_id}.json"
    path.write_text(json.dumps(data, indent=2))
    return data


def list_presentations():
    """List all saved presentations."""
    presentations = []
    for p in sorted(PRESENTATIONS_DIR.glob("*.json")):
        try:
            data = json.loads(p.read_text())
            presentations.append(data)
        except (json.JSONDecodeError, KeyError):
            continue
    return presentations


def get_presentation(pres_id):
    """Get a single presentation's metadata."""
    path = PRESENTATIONS_DIR / f"{pres_id}.json"
    if not path.exists():
        return None
    return json.loads(path.read_text())


# ───────────────────────────────────────────────────────────────────
# Gamma export (text-based)
# ───────────────────────────────────────────────────────────────────

def export_gamma_format(slides, lesson):
    """
    Export slides in a Gamma-compatible text format.
    Teachers can paste this into Gamma's "Import from text" feature.
    """
    lines = []
    lines.append(f"# {lesson.get('title', 'Lesson')}")
    lines.append(f"{lesson.get('subject', '')} — Year {lesson.get('year', '')} — Lesson {lesson.get('lesson_number', '')}")
    lines.append("")

    for slide in slides:
        title = slide.get("title", "")
        content = slide.get("content", [])
        image = slide.get("imageDescription", "")

        lines.append(f"---")
        lines.append(f"## {title}")
        for item in content:
            lines.append(f"- {item}")
        if image:
            lines.append(f"\n*Visual: {image}*")
        lines.append("")

    return "\n".join(lines)
